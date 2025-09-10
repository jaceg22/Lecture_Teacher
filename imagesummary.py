# adding images
import warnings
warnings.filterwarnings("ignore", category=UserWarning, message="resource_tracker:")

import os
# --- Set environment variables early to reduce multiprocessing issues ---
os.environ["PYTORCH_NO_CUDA_MEMORY_CACHING"] = "1"
os.environ["TOKENIZERS_PARALLELISM"] = "false"
os.environ["OMP_NUM_THREADS"] = "1"
os.environ["MKL_NUM_THREADS"] = "1"
os.environ["NUMEXPR_NUM_THREADS"] = "1"
os.environ["PYTORCH_ENABLE_MPS_FALLBACK"] = "1"  # Enable MPS fallback for Apple Silicon
os.environ["PYTORCH_MPS_HIGH_WATERMARK_RATIO"] = "0.0"  # Disable MPS memory pooling

from groq import Groq
import pytesseract
from PIL import Image
import streamlit as st
from dotenv import load_dotenv
import PyPDF2
from pptx import Presentation
from docx import Document
import io
import json
import torch
import fitz  # PyMuPDF
from docx import Document

# Use StableDiffusionPipeline for diffusers 0.34.0+
from diffusers import StableDiffusionPipeline


# Load API keys
load_dotenv()

# GROQ client for text processing (summaries, quizzes, teaching)
groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))

# Image API client configuration
# Currently using local Stable Diffusion model for image generation
# To use an external image API service instead:
# 1. Uncomment and set IMAGE_API_KEY in your .env file
# 2. Update generate_educational_image() function to use the API service
# IMAGE_API_KEY = os.getenv("IMAGE_API_KEY")  # For external image API (OpenAI DALL-E, Stability AI, etc.)

# Load image generation model with cache_resource to avoid reloads
@st.cache_resource(show_spinner=False)
def load_image_model():
    try:
        device = "cpu"  # Force CPU only to avoid Apple Silicon GPU/MPS and bitsandbytes issues
        pipe = StableDiffusionPipeline.from_pretrained(
            "runwayml/stable-diffusion-v1-5",
            torch_dtype=torch.float32,
            safety_checker=None,  # Disable safety checker to reduce memory usage
            requires_safety_checker=False
        )
        pipe = pipe.to(device)
        pipe.enable_attention_slicing()  # Reduce memory usage
        return pipe
    except Exception as e:
        st.warning(f"Failed to load image generation model: {e}")
        return None

image_pipe = load_image_model()



# Text extraction functions
def extract_text_from_image(image):
    return pytesseract.image_to_string(image)

def extract_text_from_pdf(pdf_file):
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text

def extract_text_from_pptx(pptx_file):
    presentation = Presentation(pptx_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text += shape.text + "\n"
    return text

def extract_text_from_docx(docx_file):
    doc = Document(docx_file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_images_from_pdf(pdf_file):
    """Extract images from PDF file"""
    images = []
    try:
        doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        for page_num in range(len(doc)):
            page = doc[page_num]
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image = Image.open(io.BytesIO(image_bytes))
                images.append({
                    'image': image,
                    'source': f'Page {page_num + 1}, Image {img_index + 1}',
                    'type': 'extracted'
                })
        doc.close()
    except Exception as e:
        st.warning(f"Error extracting images from PDF: {e}")
    return images

def extract_images_from_pptx(pptx_file):
    """Extract images from PowerPoint file"""
    images = []
    try:
        presentation = Presentation(pptx_file)
        for slide_num, slide in enumerate(presentation.slides):
            for shape_num, shape in enumerate(slide.shapes):
                if hasattr(shape, 'image') and shape.image:
                    image_bytes = shape.image.blob
                    image = Image.open(io.BytesIO(image_bytes))
                    images.append({
                        'image': image,
                        'source': f'Slide {slide_num + 1}, Shape {shape_num + 1}',
                        'type': 'extracted'
                    })
    except Exception as e:
        st.warning(f"Error extracting images from PowerPoint: {e}")
    return images

def extract_images_from_docx(docx_file):
    """Extract images from Word document"""
    images = []
    try:
        doc = Document(docx_file)
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_bytes = rel.target_part.blob
                image = Image.open(io.BytesIO(image_bytes))
                images.append({
                    'image': image,
                    'source': f'Word Document Image',
                    'type': 'extracted'
                })
    except Exception as e:
        st.warning(f"Error extracting images from Word document: {e}")
    return images

def extract_text_and_images(uploaded_file):
    """Extract both text and images from uploaded file"""
    file_type = uploaded_file.type
    text = ""
    images = []

    if file_type in ["image/jpeg", "image/jpg", "image/png"]:
        image = Image.open(uploaded_file)
        text = extract_text_from_image(image)
        images.append({
            'image': image,
            'source': 'Uploaded Image',
            'type': 'uploaded'
        })
    elif file_type == "application/pdf":
        text = extract_text_from_pdf(uploaded_file)
        # Reset file pointer for image extraction
        uploaded_file.seek(0)
        images = extract_images_from_pdf(uploaded_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = extract_text_from_pptx(uploaded_file)
        # Reset file pointer for image extraction
        uploaded_file.seek(0)
        images = extract_images_from_pptx(uploaded_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        text = extract_text_from_docx(uploaded_file)
        # Reset file pointer for image extraction
        uploaded_file.seek(0)
        images = extract_images_from_docx(uploaded_file)
    else:
        text = "Unsupported file type"
    
    return text, images

# AI Functions
def summarize_text(text):
    try:
        # Truncate text to avoid token limits
        truncated_text = text[:3500] if len(text) > 3500 else text
        response = groq_client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that creates detailed summaries of academic content."},
                {"role": "user", "content": f"Please create a comprehensive summary of the following content. Include key concepts, important details, formulas, examples, and any problem-solving approaches:\n{truncated_text}"}
            ],
            temperature=0.3,
            max_tokens=800
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"Error generating summary: {str(e)}\n\nContent preview:\n{text[:1000]}..."

@st.cache_data(ttl=3600)  # Cache for 1 hour
def extract_text_from_image_cached(image_bytes):
    """Cached OCR extraction to avoid re-processing same images"""
    from PIL import Image
    import io
    image = Image.open(io.BytesIO(image_bytes))
    return extract_text_from_image(image)

def analyze_images_for_content(text, extracted_images):
    """Use AI to classify images as diagrams, graphs, or tables and determine their educational value"""
    if not extracted_images:
        st.info("🔍 No images found in the uploaded document.")
        return []
    
    st.info(f"🔍 Using AI to analyze {len(extracted_images)} images...")
    import concurrent.futures
    
    def analyze_single_image(img_data, text_context):
        try:
            # Convert image to bytes for caching
            import io
            img_bytes = io.BytesIO()
            img_data['image'].save(img_bytes, format='PNG')
            img_bytes = img_bytes.getvalue()
            
            # Use cached OCR
            image_text = extract_text_from_image_cached(img_bytes)
            
            # Use AI to classify the image type and determine educational value
            response = groq_client.chat.completions.create(
                model="llama3-8b-8192",
                messages=[
                    {"role": "system", "content": "You are an expert at identifying visual educational materials. Return JSON format: {\"type\": \"diagram|graph|table|equation|text|other\", \"educational_value\": \"high|medium|low\", \"topic\": \"brief description\", \"include\": true/false, \"reason\": \"why included/excluded\"}. STRICT RULES: Only include=true for: 1) DIAGRAMS (flowcharts, process diagrams, architectural diagrams), 2) GRAPHS (charts, plots, visualizations with axes), 3) TABLES (structured data in rows/columns). EXCLUDE: equations, formulas, plain text, code snippets, simple lists."},
                    {"role": "user", "content": f"Document context: {text_context[:200]}...\n\nImage content (OCR): {image_text[:400]}...\n\nIs this a visual DIAGRAM, GRAPH, or TABLE that helps explain concepts? Be very strict - exclude anything that's just equations, formulas, or text."}
                ],
                temperature=0.1,
                max_tokens=150
            )
            
            result_text = response.choices[0].message.content.strip()
            
            # Parse JSON response
            import re
            try:
                json_match = re.search(r'\{.*\}', result_text, re.DOTALL)
                if json_match:
                    result = json.loads(json_match.group())
                else:
                    result = {"include": False, "type": "unknown", "topic": "unclassified"}
            except:
                result = {"include": False, "type": "unknown", "topic": "parsing error"}
            
            # Additional strict filtering
            visual_type = result.get("type", "").lower()
            include_decision = result.get("include", False)
            
            # Exclude equations, formulas, and text-heavy content
            if visual_type in ["equation", "text", "other"]:
                include_decision = False
            
            # Additional text-based filtering
            words = image_text.split()
            text_lower = image_text.lower()
            
            # Exclude if it's mostly mathematical equations or formulas
            equation_indicators = ['=', '+', '-', '*', '/', '^', '(', ')', 'x=', 'y=', 'f(x)', '∑', '∫', '∂']
            equation_count = sum(1 for indicator in equation_indicators if indicator in image_text)
            if equation_count > len(words) * 0.3:  # More than 30% equation symbols
                include_decision = False
            
            # Exclude if it's just a list or plain text
            if (len(words) > 20 and 
                not any(indicator in text_lower for indicator in ['table', 'chart', 'graph', 'diagram', 'flow', 'process']) and
                visual_type not in ['diagram', 'graph', 'table']):
                include_decision = False
            
            # Only include if type is specifically diagram, graph, or table AND AI said include=true
            if include_decision and visual_type in ["diagram", "graph", "table"]:
                topic = result.get("topic", "").lower()
                
                return {
                    **img_data,
                    'visual_type': visual_type,
                    'topic': topic,
                    'extracted_text': image_text[:200] + "..." if len(image_text) > 200 else image_text,
                    'educational_value': result.get("educational_value", "medium"),
                    'classification_reason': result.get("reason", "AI classified as educational visual")
                }
            
        except Exception as e:
            pass  # Skip failed images
        return None
    
    # Process images in parallel
    helpful_images = []
    with concurrent.futures.ThreadPoolExecutor(max_workers=3) as executor:
        futures = [executor.submit(analyze_single_image, img_data, text) for img_data in extracted_images]
        
        for future in concurrent.futures.as_completed(futures, timeout=30):
            try:
                result = future.result()
                if result:
                    helpful_images.append(result)
            except Exception:
                continue
    
    # Sort by educational value and type
    educational_priority = {"high": 3, "medium": 2, "low": 1}
    helpful_images.sort(key=lambda x: educational_priority.get(x.get('educational_value', 'medium'), 2), reverse=True)
    
    if helpful_images:
        types_found = [img.get('visual_type', 'visual') for img in helpful_images]
        st.success(f"✅ Found {len(helpful_images)} educational visuals: {', '.join(set(types_found))}")
        
        # Debug information
        with st.expander("🔍 Debug: Image Classification Details", expanded=False):
            for i, img in enumerate(helpful_images):
                st.write(f"**Image {i+1}:**")
                st.write(f"- Type: {img.get('visual_type', 'unknown')}")
                st.write(f"- Topic: {img.get('topic', 'unknown')}")
                st.write(f"- Educational Value: {img.get('educational_value', 'unknown')}")
                st.write(f"- Text Preview: {img.get('extracted_text', 'No text')[:100]}...")
                st.write("---")
    else:
        st.warning("⚠️ No helpful diagrams/graphs/tables found in the uploaded images.")
    
    return helpful_images  # Return all helpful images for better integration

def teach_content_with_integrated_images(text, extracted_images):
    """Generate teaching content by first analyzing images to create a structured plan"""
    
    import concurrent.futures
    
    # Step 1: Analyze and classify images to understand what visual topics we have
    helpful_images = analyze_images_for_content(text, extracted_images)
    
    if not helpful_images:
        # No helpful images - generate standard teaching content
        try:
            response = groq_client.chat.completions.create(
                model="llama3-8b-8192",
                messages=[
                    {"role": "system", "content": "You are an expert tutor creating comprehensive study materials."},
                    {"role": "user", "content": f"Create comprehensive teaching content for these study notes. For each major concept provide detailed explanations, step-by-step processes, code examples, worked examples, and comparisons.\n\nStudy notes:\n{text[:3500]}"}
                ],
                temperature=0.2,
                max_tokens=4500
            )
            teaching_content = response.choices[0].message.content.strip()
            return teaching_content, []
        except Exception as e:
            return f"Error generating teaching content: {str(e)}", []
    
    # Step 2: Create a structured teaching plan based on available images and text content
    st.info(f"📋 Creating structured teaching plan based on {len(helpful_images)} visual aids...")
    
    try:
        # Extract topics from helpful images
        image_topics = []
        for img in helpful_images:
            topic_info = {
                'topic': img.get('topic', 'unknown'),
                'visual_type': img.get('visual_type', 'diagram'),
                'extracted_text': img.get('extracted_text', ''),
                'educational_value': img.get('educational_value', 'medium')
            }
            image_topics.append(topic_info)
        
        # Create teaching plan prompt
        image_topics_text = "\n".join([
            f"- {topic['visual_type'].title()}: {topic['topic']} (text preview: {topic['extracted_text'][:100]}...)"
            for topic in image_topics
        ])
        
        # Generate structured teaching plan
        plan_response = groq_client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {"role": "system", "content": "You are an expert curriculum designer. Create a structured teaching plan that integrates visual aids with text content. Return a JSON structure with teaching sections."},
                {"role": "user", "content": f"""Create a structured teaching plan that integrates these visual aids with the text content.

AVAILABLE VISUAL AIDS:
{image_topics_text}

TEXT CONTENT TO ANALYZE:
{text[:2000]}...

Return a JSON structure like this:
{{
  "sections": [
    {{
      "title": "Section Title",
      "has_visual": true,
      "visual_topic": "exact topic from visual aids",
      "content_focus": "what to explain about this topic",
      "additional_concepts": ["other related concepts from text"]
    }}
  ],
  "additional_topics": ["important topics from text not covered by visuals"]
}}

RULES:
1. Create sections around the visual aids first
2. Identify what text content relates to each visual
3. Note additional important topics from text that don't have visuals
4. Keep visual_topic exactly matching the topics from available visual aids"""}
            ],
            temperature=0.1,
            max_tokens=800
        )
        
        # Parse the teaching plan
        plan_text = plan_response.choices[0].message.content.strip()
        import re
        import json
        
        json_match = re.search(r'\{.*\}', plan_text, re.DOTALL)
        if json_match:
            teaching_plan = json.loads(json_match.group())
        else:
            # Fallback plan if JSON parsing fails
            teaching_plan = {
                "sections": [{"title": topic['topic'], "has_visual": True, "visual_topic": topic['topic'], "content_focus": f"Explain {topic['topic']}"} for topic in image_topics],
                "additional_topics": []
            }
        
        st.success(f"✅ Created teaching plan with {len(teaching_plan.get('sections', []))} visual sections")
        
        # Debug: Show the teaching plan
        with st.expander("🔍 Debug: Teaching Plan Structure", expanded=False):
            st.json(teaching_plan)
        
    except Exception as e:
        st.warning(f"Failed to create structured plan: {e}")
        # Fallback: create simple plan based on image topics
        teaching_plan = {
            "sections": [{"title": topic['topic'], "has_visual": True, "visual_topic": topic['topic']} for topic in image_topics],
            "additional_topics": []
        }
    
    # Step 3: Generate content for each section with images integrated
    st.info("📝 Generating detailed content for each section...")
    
    teaching_sections = []
    
    # Process sections with visuals first
    for section in teaching_plan.get('sections', []):
        if section.get('has_visual', False):
            visual_topic = section.get('visual_topic', '')
            
            # Find the matching image
            matching_image = None
            for img in helpful_images:
                if visual_topic.lower() in img.get('topic', '').lower():
                    matching_image = img
                    break
            
            if matching_image:
                try:
                    # Generate focused content for this specific visual topic
                    content_response = groq_client.chat.completions.create(
                        model="llama3-8b-8192",
                        messages=[
                            {"role": "system", "content": f"You are an expert tutor explaining {visual_topic}. Write detailed educational content that will be paired with a {matching_image.get('visual_type', 'diagram')}. Focus on making the text complement the visual aid."},
                            {"role": "user", "content": f"""Write comprehensive educational content about: {visual_topic}

Context from the visual aid:
- Type: {matching_image.get('visual_type', 'diagram')}
- Content preview: {matching_image.get('extracted_text', '')[:300]}

Related text content:
{text[:1500]}

Write 2-4 paragraphs that:
1. Introduce and define the concept
2. Explain key principles and how they work
3. Provide examples or applications
4. Connect to the visual representation

Make it detailed but accessible, and write as if the reader will see a {matching_image.get('visual_type')} right after reading this."""}
                        ],
                        temperature=0.2,
                        max_tokens=600
                    )
                    
                    section_content = content_response.choices[0].message.content.strip()
                    
                    teaching_sections.append({
                        'type': 'text_with_image',
                        'title': section.get('title', visual_topic).title(),
                        'content': section_content,
                        'image': matching_image
                    })
                    
                    # Remove this image from the helpful_images list so it's not used again
                    helpful_images.remove(matching_image)
                    
                except Exception as e:
                    st.warning(f"Failed to generate content for {visual_topic}: {e}")
    
    # Step 4: Generate content for additional topics from text (no visuals)
    additional_topics = teaching_plan.get('additional_topics', [])
    if additional_topics:
        try:
            additional_content_response = groq_client.chat.completions.create(
                model="llama3-8b-8192",
                messages=[
                    {"role": "system", "content": "You are an expert tutor covering additional important concepts from the study material."},
                    {"role": "user", "content": f"""Cover these additional important topics from the study material:
{', '.join(additional_topics)}

Source content:
{text[:2000]}

Write clear explanations for each topic, including:
- Definitions and key concepts
- How they work or apply
- Examples where relevant
- Why they're important

Format with clear headings for each topic."""}
                ],
                temperature=0.2,
                max_tokens=1000
            )
            
            additional_content = additional_content_response.choices[0].message.content.strip()
            
            teaching_sections.append({
                'type': 'text_only',
                'title': 'Additional Key Concepts',
                'content': additional_content
            })
            
        except Exception as e:
            st.warning(f"Failed to generate additional topics content: {e}")
    
    # Step 5: Add any remaining helpful images with explanatory content
    if helpful_images:
        st.info(f"📊 Adding explanatory content for {len(helpful_images)} remaining visual aids...")
        
        for img in helpful_images:
            try:
                explanation_response = groq_client.chat.completions.create(
                    model="llama3-8b-8192",
                    messages=[
                        {"role": "system", "content": f"You are explaining a {img.get('visual_type', 'diagram')} about {img.get('topic', 'educational content')} to students."},
                        {"role": "user", "content": f"""Write educational content to accompany this {img.get('visual_type', 'diagram')}:

Topic: {img.get('topic', 'educational content')}
Content from image: {img.get('extracted_text', '')[:300]}

Related study material:
{text[:1500]}

Write 2-3 paragraphs that:
1. Explain what this visual demonstrates
2. Connect it to the broader learning objectives
3. Highlight key insights students should gain

Be educational and specific about how this visual aids understanding."""}
                    ],
                    temperature=0.3,
                    max_tokens=400
                )
                
                explanation = explanation_response.choices[0].message.content.strip()
                
                teaching_sections.append({
                    'type': 'text_with_image',
                    'title': f"{img.get('visual_type', 'Visual').title()}: {img.get('topic', 'Educational Content').title()}",
                    'content': explanation,
                    'image': img
                })
                
            except Exception as e:
                # Fallback explanation if AI fails
                teaching_sections.append({
                    'type': 'text_with_image',
                    'title': f"{img.get('visual_type', 'Visual').title()}: {img.get('topic', 'Educational Content').title()}",
                    'content': f"This {img.get('visual_type', 'diagram')} provides valuable insight into {img.get('topic', 'the educational content')} and complements the main teaching material by offering a visual representation of key concepts.",
                    'image': img
                })
    
    # Step 6: Combine all sections into final teaching content
    final_content = "## 🎓 Comprehensive Study Guide\n\n"
    final_content += "*This guide integrates visual learning aids with detailed explanations to help you master the concepts.*\n\n"
    
    all_images_used = []
    
    for i, section in enumerate(teaching_sections):
        final_content += f"### {section['title']}\n\n"
        final_content += section['content'] + "\n\n"
        
        if section['type'] == 'text_with_image':
            final_content += f"[PLACE_IMAGE_{i}]\n\n"
            all_images_used.append(section['image'])
        
        final_content += "---\n\n"
    
    st.success(f"✅ Generated comprehensive teaching content with {len(all_images_used)} integrated visual aids")
    
    return final_content, all_images_used


def render_integrated_teaching_content(teaching_text, images):
    """Render teaching content with images integrated at designated spots"""
    st.write("## 🎓 Exam Preparation")
    
    if not teaching_text or len(teaching_text.strip()) == 0:
        st.error("❌ No teaching content was generated. This might be due to an API issue.")
        st.info("Please try clicking the 'Teach' button again, or check if your GROQ_API_KEY is working properly.")
        return
    
    import re
    
    # Find all image placement markers in the text
    image_markers = re.findall(r'\[PLACE_IMAGE_(\d+)\]', teaching_text)
    
    # Split content by image markers
    parts = re.split(r'\[PLACE_IMAGE_\d+\]', teaching_text)
    
    # Render content with images
    for i, part in enumerate(parts):
        # Display text content
        if part.strip():
            st.write(part)
        
        # Display image if there's a corresponding marker
        if i < len(image_markers) and i < len(images):
            img_data = images[i]
            
            # Create informative caption
            visual_type = img_data.get('visual_type', 'visual').title()
            topic = img_data.get('topic', 'educational content')
            source_file = img_data.get('source_file', 'document')
            
            caption = f"📊 {visual_type}: {topic} (from {source_file})"
            
            st.image(img_data['image'], caption=caption, use_container_width=True)
            st.write("")  # Add spacing after image

def create_multiple_choice(text):
    try:
        # Truncate text to avoid token limits
        truncated_text = text[:3000] if len(text) > 3000 else text
        response = groq_client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {"role": "system", "content": "Create challenging multiple choice questions for exam preparation. You must return ONLY valid JSON format with no additional text before or after."},
                {"role": "user", "content": f"Create 6 challenging multiple choice questions from this content. Return ONLY the JSON array with no additional text:\n[{{\n  \"question\": \"What is...?\",\n  \"options\": {{\n    \"A\": \"Option A text\",\n    \"B\": \"Option B text\",\n    \"C\": \"Option C text\",\n    \"D\": \"Option D text\"\n  }},\n  \"correct_answer\": \"A\",\n  \"explanation\": \"Explanation text\"\n}}]\n\nContent to analyze:\n{truncated_text}"}
            ],
            temperature=0.4,
            max_tokens=1500
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"[]"

def create_true_false(text):
    try:
        # Truncate text to avoid token limits
        truncated_text = text[:3000] if len(text) > 3000 else text
        response = groq_client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {"role": "system", "content": "Create challenging true/false questions for exam preparation. Return ONLY a JSON array with no additional text before or after."},
                {"role": "user", "content": f"Create 8-10 challenging true/false questions from this content. Return ONLY the JSON array with format: [{{\"statement\": \"...\", \"correct_answer\": true, \"explanation\": \"...\"}}]. Make them difficult and focus on nuanced details:\n{truncated_text}"}
            ],
            temperature=0.4,
            max_tokens=1000
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"[]"

# Streamlit UI
st.title("📚 StudyLM - Your AI Study Assistant")
st.write("Upload your lecture notes, slides, and documents to generate summaries, get personalized teaching, and create practice quizzes.")

uploaded_file = st.file_uploader("📁 Upload your source", type=["jpg", "jpeg", "png", "pdf", "pptx", "docx"], accept_multiple_files=False, help="Upload lecture slides, PDFs, Word docs, or images")

if uploaded_file:
    st.subheader("📚 Your Lecture Notes")
    
    # Show file info
    file_type = uploaded_file.type
    file_size = f"{uploaded_file.size / 1024:.1f} KB" if uploaded_file.size < 1024*1024 else f"{uploaded_file.size / (1024*1024):.1f} MB"
    
    if file_type == "application/pdf":
        icon = "📄"
        type_label = "PDF"
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        icon = "📊"
        type_label = "PowerPoint"
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        icon = "📝"
        type_label = "Word"
    elif file_type in ["image/jpeg", "image/jpg", "image/png"]:
        icon = "🖼️"
        type_label = "Image"
    else:
        icon = "📁"
        type_label = "File"
    
    st.write(f"**File:** {icon} {uploaded_file.name} ({type_label}, {file_size})")
    
    # Process the single file
    with st.spinner("Processing file..."):
        text, extracted_images = extract_text_and_images(uploaded_file)
        
        # Add source file info to images
        for img_data in extracted_images:
            img_data['source_file'] = uploaded_file.name
    
    # Show preview of text
    with st.expander("📖 Preview content", expanded=False):
        st.text_area("Extracted Text", text, height=300)
    
    # Show count of extracted images
    if extracted_images:
        st.info(f"📷 Found {len(extracted_images)} image(s). Helpful images will be shown in the teaching content.")
    
    # Initialize session state
    if 'content_type' not in st.session_state:
        st.session_state.content_type = None
    if 'content_data' not in st.session_state:
        st.session_state.content_data = None
    if 'quiz_answers' not in st.session_state:
        st.session_state.quiz_answers = {}
    if 'quiz_submitted' not in st.session_state:
        st.session_state.quiz_submitted = False

    # Create 4 columns for buttons
    col1, col2, col3, col4 = st.columns(4)

    with col1:
        if st.button("📝 Generate Summary"):
            # Check if we already have cached summary for this text
            text_hash = hash(text)
            cache_key = f"summary_{text_hash}"
            
            if cache_key in st.session_state:
                summary = st.session_state[cache_key]
            else:
                with st.spinner("Summarizing..."):
                    summary = summarize_text(text)
                    # Cache the result
                    st.session_state[cache_key] = summary
            
            st.session_state.content_type = "summary"
            st.session_state.content_data = summary
            st.session_state.quiz_submitted = False

    with col2:
        if st.button("🎓 Teach"):
            # Check if we already have cached teaching content for this text
            text_hash = hash(text)
            cache_key = f"teach_improved_{text_hash}"
            cache_key_images = f"teach_images_improved_{text_hash}"
            
            if cache_key in st.session_state and cache_key_images in st.session_state:
                lesson = st.session_state[cache_key]
                all_images = st.session_state[cache_key_images]
            else:
                with st.spinner("Analyzing content and preparing lesson with integrated visuals..."):
                    lesson, all_images = teach_content_with_integrated_images(text, extracted_images)
                    # Cache the results
                    st.session_state[cache_key] = lesson
                    st.session_state[cache_key_images] = all_images
            
            st.session_state.content_type = "teach"
            st.session_state.content_data = lesson
            st.session_state.teaching_images = all_images
            st.session_state.quiz_submitted = False

    with col3:
        if st.button("📊 Multiple Choice"):
            # Check if we already have cached quiz for this text
            text_hash = hash(text)
            cache_key = f"mc_quiz_{text_hash}"
            
            if cache_key in st.session_state:
                mc_quiz = st.session_state[cache_key]
            else:
                with st.spinner("Creating quiz..."):
                    mc_quiz_raw = create_multiple_choice(text)
                    try:
                        # More robust JSON extraction
                        import re
                        # Remove any text before the first [ and after the last ]
                        json_match = re.search(r'\[.*\]', mc_quiz_raw, re.DOTALL)
                        if json_match:
                            clean_json = json_match.group(0)
                            mc_quiz = json.loads(clean_json)
                            # Cache the result
                            st.session_state[cache_key] = mc_quiz
                        else:
                            mc_quiz = []
                    except Exception as e:
                        st.error(f"Error parsing quiz: {str(e)}")
                        st.write("Raw response:", mc_quiz_raw)
                        mc_quiz = []
            
            st.session_state.content_type = "multiple_choice"
            st.session_state.content_data = mc_quiz
            st.session_state.quiz_answers = {}
            st.session_state.quiz_submitted = False

    with col4:
        if st.button("✅ True/False"):
            # Check if we already have cached quiz for this text
            text_hash = hash(text)
            cache_key = f"tf_quiz_{text_hash}"
            
            if cache_key in st.session_state:
                tf_quiz = st.session_state[cache_key]
            else:
                with st.spinner("Creating true/false quiz..."):
                    tf_quiz_raw = create_true_false(text)
                    try:
                        # More robust JSON extraction
                        import re
                        json_match = re.search(r'\[.*\]', tf_quiz_raw, re.DOTALL)
                        if json_match:
                            clean_json = json_match.group(0)
                            tf_quiz = json.loads(clean_json)
                            # Cache the result
                            st.session_state[cache_key] = tf_quiz
                        else:
                            tf_quiz = []
                    except Exception as e:
                        st.warning(f"Error parsing true/false quiz: {str(e)}")
                        tf_quiz = []
            
            st.session_state.content_type = "true_false"
            st.session_state.content_data = tf_quiz
            st.session_state.quiz_answers = {}
            st.session_state.quiz_submitted = False

    # Display content below buttons using full width
    st.markdown("---")

    if st.session_state.content_type == "summary":
        st.subheader("📌 Summary:")
        st.write(st.session_state.content_data)
    
    elif st.session_state.content_type == "teach":
        render_integrated_teaching_content(st.session_state.content_data, st.session_state.teaching_images)

    elif st.session_state.content_type == "multiple_choice" and st.session_state.content_data:
        st.subheader("📊 Multiple Choice Quiz:")
        if not st.session_state.quiz_submitted:
            for i, question in enumerate(st.session_state.content_data):
                st.write(f"**Question {i+1}:** {question['question']}")
                answer = st.radio(
                    f"Select your answer for Question {i+1}:",
                    options=list(question['options'].keys()),
                    format_func=lambda x: f"{x}: {question['options'][x]}",
                    key=f"mc_q{i}",
                    index=None
                )
                if answer:
                    st.session_state.quiz_answers[i] = answer
                st.write("")
            if st.button("Submit Quiz", type="primary"):
                st.session_state.quiz_submitted = True
                st.rerun()
        else:
            correct_count = 0
            for i, question in enumerate(st.session_state.content_data):
                user_answer = st.session_state.quiz_answers.get(i, "")
                correct_answer = question['correct_answer']
                is_correct = user_answer == correct_answer
                st.write(f"**Question {i+1}:** {question['question']}")
                for option, text in question['options'].items():
                    st.write(f"{option}: {text}")
                st.write("")
                if is_correct:
                    correct_count += 1
                    st.success(f"✅ Correct! Your answer: {user_answer}")
                else:
                    st.error(f"❌ Incorrect")
                    st.write(f"Your answer: {user_answer}")
                    st.write(f"Correct answer: {correct_answer}")
                    st.write(f"**Explanation:** {question['explanation']}")
                st.write("---")
            st.write(f"**Final Score: {correct_count}/{len(st.session_state.content_data)} ({correct_count/len(st.session_state.content_data)*100:.1f}%)**")

    elif st.session_state.content_type == "true_false" and st.session_state.content_data:
        st.subheader("✅ True/False Quiz:")
        if not st.session_state.quiz_submitted:
            for i, question in enumerate(st.session_state.content_data):
                st.write(f"**Statement {i+1}:** {question['statement']}")
                answer = st.radio(
                    f"Select your answer for Statement {i+1}:",
                    options=["True", "False"],
                    key=f"tf_q{i}",
                    index=None
                )
                if answer:
                    st.session_state.quiz_answers[i] = answer.lower() == "true"
                st.write("")
            if st.button("Submit Quiz", type="primary"):
                st.session_state.quiz_submitted = True
                st.rerun()
        else:
            correct_count = 0
            for i, question in enumerate(st.session_state.content_data):
                user_answer = st.session_state.quiz_answers.get(i, None)
                correct_answer = question['correct_answer']
                is_correct = user_answer == correct_answer
                st.write(f"**Statement {i+1}:** {question['statement']}")
                st.write("")
                if is_correct:
                    correct_count += 1
                    st.success(f"✅ Correct! Your answer: {str(user_answer)}")
                else:
                    st.error(f"❌ Incorrect")
                    st.write(f"Your answer: {str(user_answer)}")
                    st.write(f"Correct answer: {str(correct_answer)}")
                    st.write(f"**Explanation:** {question['explanation']}")
                st.write("---")
            st.write(f"**Final Score: {correct_count}/{len(st.session_state.content_data)} ({correct_count/len(st.session_state.content_data)*100:.1f}%)**")
